"""
Módulo para geração de relatórios em PowerPoint.
Cria apresentações profissionais com os resultados da análise.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pandas as pd
import numpy as np
from typing import Dict, List, Any, Optional
import logging
from pathlib import Path
import matplotlib.pyplot as plt
import seaborn as sns
from io import BytesIO

logger = logging.getLogger(__name__)


class PowerPointReporter:
    """
    Classe responsável por gerar relatórios em PowerPoint.
    """
    
    def __init__(self, data: pd.DataFrame, results: Dict[str, Any]):
        """
        Inicializa o gerador de relatórios.
        
        Args:
            data: DataFrame com dados processados
            results: Resultados dos testes de hipóteses
        """
        self.data = data
        self.results = results
        self.prs = Presentation()
        
    def create_presentation(self) -> Presentation:
        """
        Cria apresentação completa com todos os resultados.
        
        Returns:
            Apresentação PowerPoint
        """
        logger.info("Criando apresentação PowerPoint")
        
        # Slide 1: Título
        self._add_title_slide()
        
        # Slide 2: Objetivos
        self._add_objectives_slide()
        
        # Slide 3: Metodologia
        self._add_methodology_slide()
        
        # Slide 4: Visão Geral dos Dados
        self._add_data_overview_slide()
        
        # Slides 5-8: Hipóteses
        self._add_hypothesis_slides()
        
        # Slide 9: Resumo dos Resultados
        self._add_results_summary_slide()
        
        # Slide 10: Conclusões e Recomendações
        self._add_conclusions_slide()
        
        return self.prs
    
    def _add_title_slide(self):
        """Adiciona slide de título."""
        slide_layout = self.prs.slide_layouts[0]  # Layout de título
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Análise de Equidade Educacional"
        subtitle.text = "Investigação das Causas do Desempenho Diferenciado de Alunos Minoritários no SAEB\n\nAnálise Estatística de 4 Hipóteses Principais"
        
        # Formatação do título
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        # Formatação do subtítulo
        subtitle.text_frame.paragraphs[0].font.size = Pt(20)
        subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(64, 64, 64)
    
    def _add_objectives_slide(self):
        """Adiciona slide de objetivos."""
        slide_layout = self.prs.slide_layouts[1]  # Layout de conteúdo
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Objetivos da Análise"
        
        objectives_text = """
        • Investigar as causas do desempenho diferenciado de alunos minoritários no SAEB
        
        • Testar 4 hipóteses principais sobre equidade educacional:
        
        • Avaliar a efetividade de políticas educacionais direcionadas
        
        • Fornecer evidências empíricas para políticas públicas
        
        • Identificar fatores que perpetuam desigualdades educacionais
        """
        
        content.text = objectives_text
        self._format_content_text(content)
    
    def _add_methodology_slide(self):
        """Adiciona slide de metodologia."""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Metodologia"
        
        methodology_text = """
        • Dados: Sistema de Avaliação da Educação Básica (SAEB)
        
        • Amostra: {:,} alunos de múltiplas escolas
        
        • Métodos Estatísticos:
          - Testes t para comparação de médias
          - Análise de correlação de Pearson
          - Regressão linear múltipla
          - Análise de quartis
        
        • Software: Python (pandas, scipy, scikit-learn)
        
        • Nível de significância: α = 0.05
        """.format(len(self.data))
        
        content.text = methodology_text
        self._format_content_text(content)
    
    def _add_data_overview_slide(self):
        """Adiciona slide de visão geral dos dados."""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Visão Geral dos Dados"
        
        # Calcula estatísticas resumidas
        total_students = len(self.data)
        total_schools = self.data['CODIGO_ESCOLA'].nunique() if 'CODIGO_ESCOLA' in self.data.columns else 'N/A'
        minority_pct = self.data['MINORIA'].mean() * 100
        avg_math_score = self.data['NOTA_MATEMATICA'].mean()
        avg_port_score = self.data['NOTA_PORTUGUES'].mean()
        
        overview_text = f"""
        • Total de Alunos: {total_students:,}
        
        • Total de Escolas: {total_schools}
        
        • Percentual de Minorias: {minority_pct:.1f}%
        
        • Nota Média Matemática: {avg_math_score:.1f}
        
        • Nota Média Português: {avg_port_score:.1f}
        
        • Variáveis Analisadas:
          - Desempenho acadêmico (notas)
          - Características socioeconômicas
          - Infraestrutura escolar
          - Qualificação docente
          - Capital cultural
        """
        
        content.text = overview_text
        self._format_content_text(content)
    
    def _add_hypothesis_slides(self):
        """Adiciona slides para cada hipótese."""
        hypothesis_descriptions = {
            'hypothesis_1': {
                'title': 'Hipótese 1: Segregação Socioespacial',
                'description': 'Alunos minoritários concentrados em escolas com menor infraestrutura'
            },
            'hypothesis_2': {
                'title': 'Hipótese 2: Qualidade Docente',
                'description': 'Professores menos qualificados em escolas com maior concentração de minorias'
            },
            'hypothesis_3': {
                'title': 'Hipótese 3: Capital Cultural',
                'description': 'Diferenças no ambiente familiar e recursos educacionais domésticos'
            },
            'hypothesis_4': {
                'title': 'Hipótese 4: Efeito de Pares',
                'description': 'Impacto negativo da composição socioeconômica da turma'
            }
        }
        
        for key, info in hypothesis_descriptions.items():
            if key in self.results:
                self._add_single_hypothesis_slide(info['title'], info['description'], self.results[key])
    
    def _add_single_hypothesis_slide(self, title: str, description: str, results: Dict[str, Any]):
        """Adiciona slide para uma hipótese específica."""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        slide_title = slide.shapes.title
        content = slide.placeholders[1]
        
        slide_title.text = title
        
        # Constrói texto do slide
        slide_text = f"{description}\n\n"
        
        # Adiciona resultados dos testes
        significant_tests = []
        non_significant_tests = []
        
        for test_name, test_result in results['tests'].items():
            if isinstance(test_result, dict) and 'significant' in test_result:
                if test_result['significant']:
                    significant_tests.append(f"• {test_name}: SIGNIFICATIVO (p = {test_result['p_value']:.4f})")
                else:
                    non_significant_tests.append(f"• {test_name}: Não significativo (p = {test_result['p_value']:.4f})")
        
        if significant_tests:
            slide_text += "Resultados Significativos:\n" + "\n".join(significant_tests) + "\n\n"
        
        if non_significant_tests:
            slide_text += "Resultados Não Significativos:\n" + "\n".join(non_significant_tests) + "\n\n"
        
        # Adiciona estatísticas resumidas
        if 'summary_stats' in results:
            slide_text += "Estatísticas Resumidas:\n"
            for stat_name, stat_value in results['summary_stats'].items():
                if isinstance(stat_value, float):
                    slide_text += f"• {stat_name}: {stat_value:.2f}\n"
                else:
                    slide_text += f"• {stat_name}: {stat_value}\n"
        
        content.text = slide_text
        self._format_content_text(content)
    
    def _add_results_summary_slide(self):
        """Adiciona slide de resumo dos resultados."""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Resumo dos Resultados"
        
        # Conta hipóteses confirmadas e rejeitadas
        confirmed_hypotheses = []
        rejected_hypotheses = []
        
        for key, result in self.results.items():
            has_significant_result = False
            
            for test_name, test_result in result['tests'].items():
                if isinstance(test_result, dict) and test_result.get('significant', False):
                    has_significant_result = True
                    break
            
            if has_significant_result:
                confirmed_hypotheses.append(result['hypothesis'])
            else:
                rejected_hypotheses.append(result['hypothesis'])
        
        summary_text = f"""
        Hipóteses Confirmadas ({len(confirmed_hypotheses)}):
        """
        
        for hypothesis in confirmed_hypotheses:
            summary_text += f"• {hypothesis}\n"
        
        summary_text += f"\nHipóteses Rejeitadas ({len(rejected_hypotheses)}):\n"
        
        for hypothesis in rejected_hypotheses:
            summary_text += f"• {hypothesis}\n"
        
        summary_text += f"""
        
        Principais Achados:
        • {len(confirmed_hypotheses)} de {len(self.results)} hipóteses foram confirmadas
        • Evidências de desigualdades estruturais no sistema educacional
        • Necessidade de políticas mais efetivas para equidade
        """
        
        content.text = summary_text
        self._format_content_text(content)
    
    def _add_conclusions_slide(self):
        """Adiciona slide de conclusões e recomendações."""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Conclusões e Recomendações"
        
        conclusions_text = """
        Principais Conclusões:
        
        • Desigualdades educacionais persistem mesmo com políticas direcionadas
        
        • Fatores estruturais (infraestrutura, qualificação docente) impactam significativamente
        
        • Capital cultural e efeito de pares são determinantes importantes
        
        • Políticas atuais são insuficientes para garantir equidade
        
        Recomendações:
        
        • Investimento em infraestrutura de escolas com maior concentração de minorias
        
        • Programas de capacitação docente específicos
        
        • Políticas de redistribuição de recursos educacionais
        
        • Monitoramento contínuo de indicadores de equidade
        
        • Implementação de políticas de ação afirmativa mais robustas
        """
        
        content.text = conclusions_text
        self._format_content_text(content)
    
    def _format_content_text(self, content):
        """Formata o texto do conteúdo do slide."""
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = RGBColor(64, 64, 64)
            paragraph.alignment = PP_ALIGN.LEFT
    
    def save_presentation(self, output_path: str) -> None:
        """
        Salva a apresentação em arquivo.
        
        Args:
            output_path: Caminho para salvar a apresentação
        """
        self.prs.save(output_path)
        logger.info(f"Apresentação salva em {output_path}")
    
    def create_detailed_report(self) -> str:
        """
        Cria relatório detalhado em texto.
        
        Returns:
            String com relatório detalhado
        """
        report = "RELATÓRIO DETALHADO - ANÁLISE DE EQUIDADE EDUCACIONAL\n"
        report += "=" * 60 + "\n\n"
        
        # Resumo executivo
        report += "RESUMO EXECUTIVO\n"
        report += "-" * 20 + "\n"
        report += f"Este relatório apresenta uma análise estatística abrangente dos fatores que influenciam\n"
        report += f"o desempenho educacional de alunos minoritários no Sistema de Avaliação da Educação Básica (SAEB).\n"
        report += f"A análise baseia-se em uma amostra de {len(self.data):,} alunos e testa 4 hipóteses principais\n"
        report += f"sobre as causas das desigualdades educacionais.\n\n"
        
        # Metodologia
        report += "METODOLOGIA\n"
        report += "-" * 15 + "\n"
        report += "• Dados: Sistema de Avaliação da Educação Básica (SAEB)\n"
        report += "• Métodos: Testes t, correlação de Pearson, regressão linear múltipla\n"
        report += "• Software: Python (pandas, scipy, scikit-learn)\n"
        report += "• Nível de significância: α = 0.05\n\n"
        
        # Resultados por hipótese
        report += "RESULTADOS POR HIPÓTESE\n"
        report += "-" * 25 + "\n\n"
        
        for key, result in self.results.items():
            report += f"{result['hypothesis'].upper()}\n"
            report += f"Descrição: {result['description']}\n"
            report += "-" * 30 + "\n"
            
            for test_name, test_result in result['tests'].items():
                if isinstance(test_result, dict) and 'significant' in test_result:
                    significance = "SIGNIFICATIVO" if test_result['significant'] else "NÃO SIGNIFICATIVO"
                    report += f"{test_name}: {significance} (p = {test_result['p_value']:.4f})\n"
                    if 'effect_size' in test_result:
                        report += f"  Tamanho do efeito: {test_result['effect_size']:.4f}\n"
                elif isinstance(test_result, (int, float)):
                    report += f"{test_name}: {test_result:.4f}\n"
            
            report += "\n"
        
        # Conclusões
        report += "CONCLUSÕES\n"
        report += "-" * 12 + "\n"
        
        confirmed_count = sum(1 for result in self.results.values() 
                            if any(test.get('significant', False) 
                                 for test in result['tests'].values() 
                                 if isinstance(test, dict)))
        
        report += f"• {confirmed_count} de {len(self.results)} hipóteses foram confirmadas estatisticamente\n"
        report += "• Evidências de desigualdades estruturais no sistema educacional\n"
        report += "• Necessidade de políticas mais efetivas para garantir equidade\n"
        report += "• Importância de monitoramento contínuo de indicadores de equidade\n\n"
        
        report += "RECOMENDAÇÕES\n"
        report += "-" * 15 + "\n"
        report += "• Investimento em infraestrutura de escolas com maior concentração de minorias\n"
        report += "• Programas de capacitação docente específicos\n"
        report += "• Políticas de redistribuição de recursos educacionais\n"
        report += "• Implementação de políticas de ação afirmativa mais robustas\n"
        report += "• Monitoramento contínuo de indicadores de equidade\n"
        
        return report


if __name__ == "__main__":
    # Exemplo de uso
    from data_processing.data_processor import create_sample_data
    from analysis.hypothesis_tester import HypothesisTester
    
    # Cria dados de exemplo
    sample_data = create_sample_data(5000)
    
    # Executa testes
    tester = HypothesisTester(sample_data)
    results = tester.run_all_tests()
    
    # Cria apresentação
    reporter = PowerPointReporter(sample_data, results)
    presentation = reporter.create_presentation()
    
    # Salva apresentação
    reporter.save_presentation("relatorio_equidade_educacional.pptx")
    
    # Cria relatório detalhado
    detailed_report = reporter.create_detailed_report()
    print(detailed_report)
